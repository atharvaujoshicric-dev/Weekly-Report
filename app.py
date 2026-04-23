import streamlit as st
import pandas as pd
from pptx import Presentation
import io

st.set_page_config(page_title="BeyondWalls Report Generator", layout="wide")

st.title("📊 BeyondWalls PPT Report Generator")

# 1. Configuration & Template Upload
st.sidebar.header("Template Setup")
template_file = st.sidebar.file_uploader("Upload your PPTX Template", type="pptx")

report_type = st.sidebar.selectbox(
    "Select Report Type",
    ["CP Aggregator Weekly", "Pre-Sales Review", "Weekly Project Review"]
)

if template_file:
    st.success("Template Loaded!")
    
    # 2. Data Entry UI
    with st.form("data_entry"):
        st.subheader(f"Enter Data for {report_type}")
        col1, col2 = st.columns(2)
        
        with col1:
            proj = st.text_input("Project Name", value="Aishwaryam Abhimaan") [cite: 4]
            leads = st.number_input("Total Leads", value=0) [cite: 237]
            spends = st.number_input("Total Spends (with GST)", value=0) [cite: 11]
        
        with col2:
            svc = st.number_input("Site Visits Conducted", value=0) [cite: 242]
            bookings = st.number_input("Total Bookings", value=0) [cite: 338]
            date_range = st.text_input("Date Range", value="1st - 7th Sept") [cite: 5]

        submit = st.form_submit_button("Generate PPTX")

    # 3. PPTX Manipulation Logic
    if submit:
        # Load the presentation from the uploaded template
        prs = Presentation(template_file)
        
        # Simple search-and-replace logic for text placeholders
        # Note: For images, you would use prs.slides[index].shapes.add_picture()
        slides_updated = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # Replace keys with your inputs
                    if "{{PROJECT_NAME}}" in shape.text:
                        shape.text = shape.text.replace("{{PROJECT_NAME}}", proj)
                    if "{{TOTAL_LEADS}}" in shape.text:
                        shape.text = shape.text.replace("{{TOTAL_LEADS}}", str(leads))
                    if "{{TOTAL_SPENDS}}" in shape.text:
                        shape.text = shape.text.replace("{{TOTAL_SPENDS}}", str(spends))
                    if "{{DATE_RANGE}}" in shape.text:
                        shape.text = shape.text.replace("{{DATE_RANGE}}", date_range)
        
        # Save to a byte buffer for download
        binary_output = io.BytesIO()
        prs.save(binary_output)
        
        st.download_button(
            label="📩 Download Processed PPT Report",
            data=binary_output.getvalue(),
            file_name=f"{proj}_Weekly_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
else:
    st.info("Please upload one of the BeyondWalls PPTX templates to begin.") [cite: 1, 234, 282]
